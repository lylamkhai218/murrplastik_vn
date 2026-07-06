import os
import sys
import subprocess

def install_package(package, pip_name=None):
    if pip_name is None:
        pip_name = package
    try:
        __import__(package)
    except ImportError:
        print(f"Installing {pip_name}...")
        subprocess.check_call([sys.executable, "-m", "pip", "install", pip_name])

# Make sure python-pptx and Pillow are installed
install_package("pptx", "python-pptx")
install_package("PIL", "Pillow")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

def main():
    prs = Presentation()
    
    # Set slide dimensions to widescreen 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Define Color Palette (Clinical Sea Blue Theme)
    c_navy = RGBColor(11, 30, 54)        # #0B1E36 - Deep Navy-Slate
    c_blue = RGBColor(0, 119, 182)       # #0077B6 - Clinical Sea Blue
    c_cyan = RGBColor(6, 182, 212)       # #06B6D4 - Active Cyan (Accent)
    c_grey_bg = RGBColor(244, 248, 250)  # #F4F8FA - Sterile Background
    c_charcoal = RGBColor(71, 85, 105)   # #475569 - Slate-600 Body Text
    c_white = RGBColor(255, 255, 255)    # #FFFFFF - Card Background
    c_border = RGBColor(226, 237, 243)   # #E2EDF3 - Soft blue-tint border
    
    blank_slide_layout = prs.slide_layouts[6]
    
    # Image Directories
    root_dir = r"d:\T&TVina\Campain\F&B"
    prod_img_dir = os.path.join(root_dir, "Ảnh sản phẩm")
    real_img_dir = os.path.join(root_dir, "Hình ảnh thực tế")
    
    # Helper: Set slide background color
    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color
        
    # Helper: Add title with custom line separator (Catalog Style)
    def add_slide_header(slide, number, title_text):
        # Number block
        num_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(0.8), Inches(0.5))
        tf_num = num_box.text_frame
        tf_num.word_wrap = True
        p_num = tf_num.paragraphs[0]
        p_num.text = number
        p_num.font.name = 'Plus Jakarta Sans'
        p_num.font.size = Pt(28)
        p_num.font.bold = True
        p_num.font.color.rgb = c_cyan
        
        # Title text block
        title_box = slide.shapes.add_textbox(Inches(1.4), Inches(0.4), Inches(11.3), Inches(0.5))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text.upper()
        p_title.font.name = 'Plus Jakarta Sans'
        p_title.font.size = Pt(20)
        p_title.font.bold = True
        p_title.font.color.rgb = c_navy
        
        # Thin divider line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.0), Inches(12.13), Inches(0.01)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = c_border
        line.line.fill.background()
        
        # Accent indicator
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.98), Inches(1.5), Inches(0.03)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = c_cyan
        accent_bar.line.fill.background()
        
    # Helper: Add card container
    def add_card(slide, left, top, width, height, fill_color=c_white, border_color=c_border):
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = fill_color
        card.line.color.rgb = border_color
        card.line.width = Pt(0.75)
        try:
            card.adjustments[0] = 0.04  # Subtle rounded corners
        except:
            pass
        return card
        
    # Helper: Add text inside a rectangle with custom styling
    def add_textbox_styled(slide, left, top, width, height, text_list, align=PP_ALIGN.LEFT):
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.1)
        tf.margin_bottom = Inches(0.1)
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        
        for idx, item in enumerate(text_list):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = item.get("text", "")
            p.alignment = align
            p.space_after = Pt(item.get("space_after", 4))
            
            font = p.font
            requested_font = item.get("font_name", "Calibri")
            if requested_font == "Trebuchet MS":
                font.name = "Plus Jakarta Sans"
            elif requested_font == "Calibri":
                font.name = "Inter"
            else:
                font.name = requested_font
            font.size = Pt(item.get("size", 11))
            font.bold = item.get("bold", False)
            font.italic = item.get("italic", False)
            font.color.rgb = item.get("color", c_charcoal)
            
        return textbox

    # Helper: Add fit picture to slide to avoid distortion
    def add_fit_picture_styled(slide, img_path, left, top, max_width, max_height):
        if not os.path.exists(img_path):
            print(f"Warning: Image path not found: {img_path}")
            # Draw a placeholder box
            placeholder = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, max_width, max_height)
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = RGBColor(220, 220, 220)
            placeholder.line.color.rgb = c_border
            tx = placeholder.text_frame
            p = tx.paragraphs[0]
            p.text = f"Ảnh sản phẩm\n{os.path.basename(img_path)}"
            p.font.name = 'Inter'
            p.font.size = Pt(10)
            p.font.color.rgb = c_charcoal
            p.alignment = PP_ALIGN.CENTER
            try:
                placeholder.adjustments[0] = 0.04
            except:
                pass
            return placeholder
            
        try:
            with Image.open(img_path) as img:
                img_w, img_h = img.size
            aspect = img_w / img_h
            box_aspect = max_width / max_height
            
            if aspect > box_aspect:
                width = max_width
                height = max_width / aspect
            else:
                height = max_height
                width = max_height * aspect
                
            final_left = left + (max_width - width) / 2
            final_top = top + (max_height - height) / 2
            
            pic = slide.shapes.add_picture(img_path, final_left, final_top, width, height)
            return pic
        except Exception as e:
            print(f"Error adding picture {img_path}: {e}")
            return None

    # Helper: Add stylized table
    def add_styled_table(slide, rows, cols, left, top, width, height, data, col_widths=None):
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        
        if col_widths:
            for c_idx, w in enumerate(col_widths):
                table.columns[c_idx].width = w
                
        for r_idx in range(rows):
            for c_idx in range(cols):
                cell = table.cell(r_idx, c_idx)
                cell.text = data[r_idx][c_idx]
                
                # Format cell background
                cell.fill.solid()
                if r_idx == 0:
                    cell.fill.fore_color.rgb = c_blue # Header row
                else:
                    if r_idx % 2 == 1:
                        cell.fill.fore_color.rgb = c_white
                    else:
                        cell.fill.fore_color.rgb = RGBColor(240, 244, 248) # Light blue-grey alternating row
                
                # Format cell text
                tf = cell.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT if c_idx != 0 else PP_ALIGN.CENTER
                p.font.name = 'Inter'
                p.font.size = Pt(9 if r_idx > 0 else 10)
                p.font.bold = True if r_idx == 0 or c_idx == 0 else False
                p.font.color.rgb = c_white if r_idx == 0 else c_charcoal
                
        return table_shape

    # ==========================================
    # SLIDE 1: TRANG BÌA (COVER SLIDE)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide1, c_navy)
    
    # Left block background accent
    left_accent = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(7.0), Inches(7.5)
    )
    left_accent.fill.solid()
    left_accent.fill.fore_color.rgb = RGBColor(2, 3, 70) # Darker Navy
    left_accent.line.fill.background()
    
    # Large vertical Text logo (Watermark-like styling)
    watermark_box = slide1.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(6.0), Inches(1.5))
    tf_wm = watermark_box.text_frame
    p_wm = tf_wm.paragraphs[0]
    p_wm.text = "T&T VINA  |  MURRPLASTIK"
    p_wm.font.name = 'Trebuchet MS'
    p_wm.font.size = Pt(14)
    p_wm.font.bold = True
    p_wm.font.color.rgb = c_cyan
    
    # Main Cover Titles
    title_text = [
        {"text": "GIẢI PHÁP ĐI DÂY CÁP\nHỢP VỆ SINH", "font_name": "Trebuchet MS", "size": 38, "bold": True, "color": c_white, "space_after": 10},
        {"text": "CHO NHÀ MÁY THỰC PHẨM & ĐỒ UỐNG (F&B)", "font_name": "Trebuchet MS", "size": 18, "bold": True, "color": c_cyan, "space_after": 20},
        {"text": "Đạt chuẩn FDA, EHEDG & Ecolab – An toàn sinh học tuyệt đối, chống ăn mòn hóa chất tẩy rửa mạnh.", "font_name": "Calibri", "size": 12, "color": RGBColor(200, 220, 240), "space_after": 40},
        {"text": "NHÀ PHÂN PHỐI CHÍNH THỨC TẠI VIỆT NAM\nCÔNG TY TNHH CÔNG NGHIỆP T&T VINA", "font_name": "Trebuchet MS", "size": 11, "bold": True, "color": c_white, "space_after": 2}
    ]
    add_textbox_styled(slide1, Inches(0.5), Inches(1.8), Inches(6.0), Inches(5.0), title_text)
    
    # Cover Image on the right (WINCOME style split)
    add_fit_picture_styled(slide1, os.path.join(prod_img_dir, "Multi_products.jpg"), Inches(7.3), Inches(0.5), Inches(5.5), Inches(6.5))
    
    # Overlay card for the cover image (glassmorphic border detail)
    img_border_card = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(7.3), Inches(0.5), Inches(5.5), Inches(6.5)
    )
    img_border_card.fill.background()
    img_border_card.line.color.rgb = c_cyan
    img_border_card.line.width = Pt(1.5)

    # ==========================================
    # SLIDE 2: THÁCH THỨC VẬN HÀNH
    # ==========================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide2, c_grey_bg)
    add_slide_header(slide2, "01", "Thách Thức Vận Hành Trong Nhà Máy F&B")
    
    # 3 Column Cards
    card_width = Inches(3.7)
    card_height = Inches(4.8)
    card_y = Inches(1.6)
    
    challenges = [
        {
            "num": "01",
            "title": "VỆ SINH DỊCH TỄ & AN TOÀN SINH HỌC",
            "desc": "Kiểm soát nghiêm ngặt nguy cơ nhiễm khuẩn chéo. Các điểm mù, khe hở cơ khí, ren vít thông thường trong hệ thống đi dây là nơi cực kỳ lý tưởng để bụi bẩn và vi sinh vật tích tụ."
        },
        {
            "num": "02",
            "title": "TỐI ƯU HÓA THỜI GIAN VỆ SINH (DOWNTIME)",
            "desc": "Thời gian dừng máy để lau chùi, xịt rửa hóa chất ảnh hưởng trực tiếp đến năng suất. Nhà máy luôn cần các giải pháp thiết kế hợp vệ sinh (Hygienic Design) để dễ dàng xịt rửa nhanh chóng."
        },
        {
            "num": "03",
            "title": "MÔI TRƯỜNG TẨY RỬA KHẮC NGHIỆT",
            "desc": "Thiết bị điện và cáp dẫn phải chịu đựng được tần suất xịt rửa áp lực cao hàng ngày (IP69K) cùng sự tàn phá của các hóa chất tẩy rửa có tính axit, kiềm mạnh mà không bị ăn mòn, gỉ sét."
        }
    ]
    
    for idx, chall in enumerate(challenges):
        x = Inches(0.6 + idx * 4.2)
        add_card(slide2, x, card_y, card_width, card_height)
        
        # Text content
        text_data = [
            {"text": chall["num"], "font_name": "Trebuchet MS", "size": 32, "bold": True, "color": c_cyan, "space_after": 10},
            {"text": chall["title"], "font_name": "Trebuchet MS", "size": 13, "bold": True, "color": c_navy, "space_after": 15},
            {"text": chall["desc"], "font_name": "Calibri", "size": 11, "color": c_charcoal, "space_after": 4}
        ]
        add_textbox_styled(slide2, x + Inches(0.2), card_y + Inches(0.2), card_width - Inches(0.4), card_height - Inches(0.4), text_data)

    # ==========================================
    # SLIDE 3: GIẢI PHÁP MURRPLASTIK
    # ==========================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide3, c_grey_bg)
    add_slide_header(slide3, "02", "Giải Pháp Vệ Sinh Toàn Diện Từ Murrplastik")
    
    # Left Column: Principles & Certifications
    left_x = Inches(0.6)
    left_y = Inches(1.6)
    left_w = Inches(6.0)
    
    left_text = [
        {"text": "ĐIỂM KHÁC BIỆT CỦA CÔNG NGHỆ HYGIENIC DESIGN", "font_name": "Trebuchet MS", "size": 15, "bold": True, "color": c_navy, "space_after": 15},
        {"text": "• Không góc chết, dễ lau chùi:", "font_name": "Calibri", "size": 12, "bold": True, "color": c_blue, "space_after": 2},
        {"text": "Tất cả các cạnh bo tròn với bán kính tối thiểu 3 mm. Loại bỏ hoàn toàn khe hở và điểm mù kết cấu nhằm triệt tiêu cơ hội xâm nhập của vi khuẩn.", "font_name": "Calibri", "size": 11, "color": c_charcoal, "space_after": 10},
        {"text": "• Chứng nhận tiêu chuẩn khắt khe nhất:", "font_name": "Calibri", "size": 12, "bold": True, "color": c_blue, "space_after": 2},
        {"text": "Nguyên liệu đạt chuẩn FDA (Hoa Kỳ), thiết kế cơ khí tuân thủ nghiêm ngặt EHEDG (Châu Âu), kiểm nghiệm khả năng kháng hóa chất Ecolab, đáp ứng tiêu chuẩn EC 1935/2004 và EU 10/2011.", "font_name": "Calibri", "size": 11, "color": c_charcoal, "space_after": 10},
        {"text": "• Thép không gỉ cao cấp đánh bóng:", "font_name": "Calibri", "size": 12, "bold": True, "color": c_blue, "space_after": 2},
        {"text": "Khung thép không gỉ V2A (1.4301) và V4A (1.4404) được đánh bóng cơ học đạt độ nhám bề mặt Ra < 0.8 µm, ngăn ngừa sự bám dính của các phân tử hữu cơ thực phẩm.", "font_name": "Calibri", "size": 11, "color": c_charcoal, "space_after": 4}
    ]
    add_textbox_styled(slide3, left_x, left_y, left_w, Inches(5.0), left_text)
    
    # Right Column: Product Image with Spec Lens Overlay
    right_x = Inches(7.0)
    right_y = Inches(1.6)
    right_w = Inches(5.7)
    right_h = Inches(4.5)
    
    add_card(slide3, right_x, right_y, right_w, right_h)
    add_fit_picture_styled(slide3, os.path.join(prod_img_dir, "Conduits_and_fittings.jpg"), right_x + Inches(0.2), right_y + Inches(0.2), right_w - Inches(0.4), right_h - Inches(0.4))
    
    # Spec Lens Card (Signature Overlay)
    lens_x = right_x + Inches(3.0)
    lens_y = right_y + Inches(2.8)
    lens_w = Inches(2.5)
    lens_h = Inches(1.5)
    add_card(slide3, lens_x, lens_y, lens_w, lens_h, fill_color=c_navy, border_color=c_cyan)
    
    lens_text = [
        {"text": "ĐỘ NHÁM BỀ MẶT", "font_name": "Trebuchet MS", "size": 9, "bold": True, "color": c_cyan, "space_after": 2},
        {"text": "Ra < 0.8 µm", "font_name": "Trebuchet MS", "size": 18, "bold": True, "color": c_white, "space_after": 4},
        {"text": "Vật liệu V2A / V4A đạt chuẩn vệ sinh tối ưu.", "font_name": "Calibri", "size": 9, "color": RGBColor(220, 230, 245), "space_after": 0}
    ]
    add_textbox_styled(slide3, lens_x + Inches(0.15), lens_y + Inches(0.15), lens_w - Inches(0.3), lens_h - Inches(0.3), lens_text)

    # ==========================================
    # SLIDE 4: TÙY BIẾN THEO YÊU CẦU (KDP ON DEMAND)
    # ==========================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide4, c_grey_bg)
    add_slide_header(slide4, "03", "1.2 Tấm Dẫn Cáp Tùy Biến (KDP on Demand)")
    
    # Left Column: Service details
    left_x = Inches(0.6)
    left_y = Inches(1.6)
    left_w = Inches(6.0)
    
    left_text = [
        {"text": "GIẢI PHÁP THIẾT KẾ CÁ NHÂN HÓA DUY NHẤT", "font_name": "Trebuchet MS", "size": 15, "bold": True, "color": c_navy, "space_after": 15},
        {"text": "• Tự do cấu hình kích thước & số lỗ:", "font_name": "Calibri", "size": 12, "bold": True, "color": c_blue, "space_after": 2},
        {"text": "Được sản xuất theo bản vẽ lỗ luồn cáp chính xác của khách hàng, tương thích tuyệt đối với số lượng cáp thực tế.", "font_name": "Calibri", "size": 11, "color": c_charcoal, "space_after": 12},
        {"text": "• Lắp đặt nhanh gấp 5 lần:", "font_name": "Calibri", "size": 12, "bold": True, "color": c_blue, "space_after": 2},
        {"text": "Thao tác đơn giản bằng cách luồn trực tiếp cáp qua màng chắn dẻo. Tiết kiệm không gian cực lớn so với việc khoan hàng chục lỗ để bắt ốc siết cáp (cable gland) riêng lẻ.", "font_name": "Calibri", "size": 11, "color": c_charcoal, "space_after": 12},
        {"text": "• Mật độ luồn cáp cực cao:", "font_name": "Calibri", "size": 12, "bold": True, "color": c_blue, "space_after": 2},
        {"text": "Cho phép đi lượng lớn cáp động lực, cáp tín hiệu qua một diện tích khoét lỗ nhỏ trên thành tủ điện hoặc vách máy.", "font_name": "Calibri", "size": 11, "color": c_charcoal, "space_after": 4}
    ]
    add_textbox_styled(slide4, left_x, left_y, left_w, Inches(5.0), left_text)
    
    # Right Column: KDP on Demand image and specification table
    right_x = Inches(7.0)
    right_y = Inches(1.6)
    right_w = Inches(5.7)
    right_h = Inches(3.0)
    
    add_card(slide4, right_x, right_y, right_w, right_h)
    add_fit_picture_styled(slide4, os.path.join(prod_img_dir, "KDP 24.jpg"), right_x + Inches(0.2), right_y + Inches(0.2), right_w - Inches(0.4), right_h - Inches(0.4))
    
    # Specifications Table below image
    table_data = [
        ["THÔNG SỐ", "CHI TIẾT TÙY CHỌN KDP ON DEMAND"],
        ["Vật liệu tấm", "Thép không gỉ V2A/V4A đánh bóng hoặc Nhựa trắng thực phẩm"],
        ["Vật liệu màng bọc", "EPDM đạt chuẩn FDA hoặc Silicone màu xanh lam dễ nhận biết"],
        ["Cấu hình cơ khí", "Thiết kế vát cạnh chống tích tụ nước, bu lông hàn mặt sau"],
        ["Đường kính cáp", "Tự do lựa chọn từ 3.0 mm đến 33.0 mm cho từng lỗ luồn"]
    ]
    add_styled_table(slide4, 5, 2, right_x, Inches(4.8), right_w, Inches(2.0), table_data, col_widths=[Inches(1.8), Inches(3.9)])

    # ==========================================
    # SLIDE 5: DÒNG TẤM KIM LOẠI (PREMIUM)
    # ==========================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide5, c_grey_bg)
    add_slide_header(slide5, "04", "Tấm Dẫn Cáp Thép Không Gỉ - Dòng Premium")
    
    # Left Card: KDP/S-FDA-HD
    x_kdp = Inches(0.6)
    y_cards = Inches(1.6)
    w_cards = Inches(3.8)
    h_cards = Inches(3.3)
    
    add_card(slide5, x_kdp, y_cards, w_cards, h_cards)
    kdp_text = [
        {"text": "KDP/S-FDA-HD", "font_name": "Trebuchet MS", "size": 15, "bold": True, "color": c_navy, "space_after": 10},
        {"text": "• Thiết kế tấm kim loại liền khối bằng thép không gỉ V2A hoặc V4A.", "font_name": "Calibri", "size": 11, "color": c_charcoal, "space_after": 6},
        {"text": "• Bu lông ren được hàn chắc chắn ở mặt sau, giúp mặt trước phẳng mịn hoàn toàn, triệt tiêu khe hở đọng bụi bẩn.", "font_name": "Calibri", "size": 11, "color": c_charcoal, "space_after": 6},
        {"text": "• Đạt tiêu chuẩn kín khít tuyệt đối IP69K.", "font_name": "Calibri", "size": 11, "bold": True, "color": c_blue, "space_after": 4}
    ]
    add_textbox_styled(slide5, x_kdp + Inches(0.15), y_cards + Inches(0.15), w_cards - Inches(0.3), h_cards - Inches(0.3), kdp_text)
    
    # Right Card: KDL/H-VA-FDA
    x_kdl = Inches(4.6)
    add_card(slide5, x_kdl, y_cards, w_cards, h_cards)
    kdl_text = [
        {"text": "KDL/H-VA-FDA", "font_name": "Trebuchet MS", "size": 15, "bold": True, "color": c_navy, "space_after": 10},
        {"text": "• Khung kim loại tháo lắp, chia tách làm hai nửa, làm bằng thép không gỉ V2A.", "font_name": "Calibri", "size": 11, "color": c_charcoal, "space_after": 6},
        {"text": "• Cho phép luồn và định vị trực tiếp cáp đã có sẵn phích cắm (pre-assembled) mà không cần cắt jack nối.", "font_name": "Calibri", "size": 11, "color": c_charcoal, "space_after": 6},
        {"text": "• Tích hợp giảm căng cáp, sức chứa tới 48 cáp, bảo vệ IP69K.", "font_name": "Calibri", "size": 11, "bold": True, "color": c_blue, "space_after": 4}
    ]
    add_textbox_styled(slide5, x_kdl + Inches(0.15), y_cards + Inches(0.15), w_cards - Inches(0.3), h_cards - Inches(0.3), kdl_text)
    
    # Product Image
    x_img = Inches(8.6)
    add_card(slide5, x_img, y_cards, Inches(4.1), Inches(3.3))
    add_fit_picture_styled(slide5, os.path.join(prod_img_dir, "KDH_KDL_H-VA-FDA_Produkt.png"), x_img + Inches(0.15), y_cards + Inches(0.15), Inches(3.8), Inches(3.0))
    
    # Big Number Callouts on the side
    num_box1 = slide5.shapes.add_textbox(Inches(0.6), Inches(5.0), Inches(2.0), Inches(0.8))
    p_num1 = num_box1.text_frame.paragraphs[0]
    p_num1.text = "IP69K"
    p_num1.font.name = 'Plus Jakarta Sans'
    p_num1.font.size = Pt(24)
    p_num1.font.bold = True
    p_num1.font.color.rgb = c_cyan
    
    num_box2 = slide5.shapes.add_textbox(Inches(4.6), Inches(5.0), Inches(2.0), Inches(0.8))
    p_num2 = num_box2.text_frame.paragraphs[0]
    p_num2.text = "V4A STEEL"
    p_num2.font.name = 'Plus Jakarta Sans'
    p_num2.font.size = Pt(24)
    p_num2.font.bold = True
    p_num2.font.color.rgb = c_navy
    
    # Product Comparison Table
    table_data = [
        ["MODEL SẢN PHẨM", "CHẤT LIỆU", "CẤP BẢO VỆ", "ĐẶC TÍNH LUỒN CÁP", "ỨNG DỤNG TIÊU BIỂU"],
        ["KDP/S-FDA-HD", "Thép V2A / V4A", "IP69K", "Cáp trần (chưa bấm đầu)", "Thành tủ điện lắp ở vị trí tiếp xúc trực tiếp thực phẩm"],
        ["KDL/H-VA-FDA", "Thép V2A / Gioăng TPE", "IP69K", "Cáp đã bấm sẵn đầu cắm", "Đường truyền dữ liệu máy tính, tủ điều khiển robot"]
    ]
    add_styled_table(slide5, 3, 5, Inches(0.6), Inches(5.8), Inches(12.13), Inches(1.2), table_data, col_widths=[Inches(1.8), Inches(1.5), Inches(1.2), Inches(2.8), Inches(4.83)])

    # ==========================================
    # SLIDE 6: DÒNG TẤM NHỰA THỰC PHẨM (STANDARD)
    # ==========================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide6, c_grey_bg)
    add_slide_header(slide6, "05", "Tấm Dẫn Cáp Nhựa Thực Phẩm - Dòng Standard")
    
    # Left Card: KDP/Z-FDA
    add_card(slide6, x_kdp, y_cards, w_cards, h_cards)
    kdp_z_text = [
        {"text": "KDP/Z-FDA", "font_name": "Trebuchet MS", "size": 15, "bold": True, "color": c_navy, "space_after": 10},
        {"text": "• Khung nhựa trắng nguyên khối đạt chuẩn FDA, màng cao su kết nối kép dẻo dai.", "font_name": "Calibri", "size": 11, "color": c_charcoal, "space_after": 6},
        {"text": "• Các cạnh lỗ bo tròn giúp bảo vệ vỏ cáp không bị trầy xước trong môi trường rung động động lực học.", "font_name": "Calibri", "size": 11, "color": c_charcoal, "space_after": 6},
        {"text": "• Cấp bảo vệ kín khít đạt IP69K.", "font_name": "Calibri", "size": 11, "bold": True, "color": c_blue, "space_after": 4}
    ]
    add_textbox_styled(slide6, x_kdp + Inches(0.15), y_cards + Inches(0.15), w_cards - Inches(0.3), h_cards - Inches(0.3), kdp_z_text)
    
    # Right Card: KDL/D-FDA
    add_card(slide6, x_kdl, y_cards, w_cards, h_cards)
    kdl_d_text = [
        {"text": "KDL/D-FDA", "font_name": "Trebuchet MS", "size": 15, "bold": True, "color": c_navy, "space_after": 10},
        {"text": "• Hệ thống khung bằng nhựa trắng có thể tháo rời hoàn toàn để thi công cáp có sẵn đầu nối.", "font_name": "Calibri", "size": 11, "color": c_charcoal, "space_after": 6},
        {"text": "• Khung nhựa trắng giúp người vận hành dễ dàng phát hiện các vết bẩn, nấm mốc tích tụ bằng mắt thường.", "font_name": "Calibri", "size": 11, "color": c_charcoal, "space_after": 6},
        {"text": "• Đạt tiêu chuẩn bảo vệ IP65, tích hợp giảm căng cáp.", "font_name": "Calibri", "size": 11, "bold": True, "color": c_blue, "space_after": 4}
    ]
    add_textbox_styled(slide6, x_kdl + Inches(0.15), y_cards + Inches(0.15), w_cards - Inches(0.3), h_cards - Inches(0.3), kdl_d_text)
    
    # Product Image
    add_card(slide6, x_img, y_cards, Inches(4.1), Inches(3.3))
    add_fit_picture_styled(slide6, os.path.join(prod_img_dir, "acs_1.jpg"), x_img + Inches(0.15), y_cards + Inches(0.15), Inches(3.8), Inches(3.0))
    
    # Big Number Callouts on the side
    num_box1_s6 = slide6.shapes.add_textbox(Inches(0.6), Inches(5.0), Inches(2.0), Inches(0.8))
    p_num1_s6 = num_box1_s6.text_frame.paragraphs[0]
    p_num1_s6.text = "FDA COMPLIANT"
    p_num1_s6.font.name = 'Plus Jakarta Sans'
    p_num1_s6.font.size = Pt(22)
    p_num1_s6.font.bold = True
    p_num1_s6.font.color.rgb = c_cyan
    
    num_box2_s6 = slide6.shapes.add_textbox(Inches(4.6), Inches(5.0), Inches(2.0), Inches(0.8))
    p_num2_s6 = num_box2_s6.text_frame.paragraphs[0]
    p_num2_s6.text = "EASY DETECTION"
    p_num2_s6.font.name = 'Plus Jakarta Sans'
    p_num2_s6.font.size = Pt(22)
    p_num2_s6.font.bold = True
    p_num2_s6.font.color.rgb = c_navy
    
    # Product Comparison Table
    table_data_s6 = [
        ["MODEL SẢN PHẨM", "CHẤT LIỆU KHUNG", "CẤP BẢO VỆ", "ĐẶC ĐIỂM KỸ THUẬT", "ƯU ĐIỂM VẬN HÀNH"],
        ["KDP/Z-FDA", "Nhựa Polyamide trắng FDA", "IP69K", "Luồn cáp trần, bo góc 3D bảo vệ cáp", "Lắp đặt nhanh, giá thành tối ưu cho dây chuyền đóng gói"],
        ["KDL/D-FDA", "Nhựa Polyamide trắng FDA", "IP65", "Chia đôi khung, luồn tới 48 cáp có phích", "Dễ dàng thay thế cáp lỗi mà không cần tháo dây đấu nối"]
    ]
    add_styled_table(slide6, 3, 5, Inches(0.6), Inches(5.8), Inches(12.13), Inches(1.2), table_data_s6, col_widths=[Inches(1.8), Inches(2.0), Inches(1.2), Inches(3.0), Inches(4.13)])

    # ==========================================
    # SLIDE 7: XÍCH DẪN CÁP ĐỘNG LỰC & ỐNG BẢO VỆ CÁP
    # ==========================================
    slide7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide7, c_grey_bg)
    add_slide_header(slide7, "06", "Xích Dẫn Cáp Động Lực & Ống Bảo Vệ Cáp")
    
    # Left Half: Evochain
    x_left = Inches(0.6)
    y_split = Inches(1.6)
    w_split = Inches(5.8)
    h_split = Inches(5.2)
    
    add_card(slide7, x_left, y_split, w_split, h_split)
    add_fit_picture_styled(slide7, os.path.join(prod_img_dir, "Evochain_420_konfektioniert_freisteller.png"), x_left + Inches(0.2), y_split + Inches(0.2), w_split - Inches(0.4), Inches(2.2))
    
    evochain_text = [
        {"text": "EVOCHAIN® MAX MP 420 (XÍCH DẪN ĐỘNG)", "font_name": "Trebuchet MS", "size": 13, "bold": True, "color": c_navy, "space_after": 8},
        {"text": "• Xích dẫn cáp tối ưu cho các chuyển động tịnh tiến tốc độ cao trên băng chuyền sản xuất thực phẩm.", "font_name": "Calibri", "size": 10.5, "color": c_charcoal, "space_after": 4},
        {"text": "• Thiết kế tối giản khe kẽ, chống bám dính hữu cơ và nấm mốc.", "font_name": "Calibri", "size": 10.5, "color": c_charcoal, "space_after": 4},
        {"text": "• Cấu trúc cơ học chịu uốn gập hàng triệu chu kỳ hoạt động liên tục.", "font_name": "Calibri", "size": 10.5, "color": c_charcoal, "space_after": 4}
    ]
    add_textbox_styled(slide7, x_left + Inches(0.2), y_split + Inches(2.6), w_split - Inches(0.4), Inches(2.3), evochain_text)
    
    # Right Half: Conduits
    x_right = Inches(6.93)
    add_card(slide7, x_right, y_split, w_split, h_split)
    add_fit_picture_styled(slide7, os.path.join(prod_img_dir, "murrplastik_EWL-PAB.png"), x_right + Inches(0.2), y_split + Inches(0.2), w_split - Inches(0.4), Inches(2.2))
    
    conduit_text = [
        {"text": "ỐNG BẢO VỆ CÁP & ĐẦU NỐI VỆ SINH", "font_name": "Trebuchet MS", "size": 13, "bold": True, "color": c_navy, "space_after": 8},
        {"text": "• Ống luồn dẻo EWX-PAE & EWL-PAB làm bằng nhựa cao cấp chịu hóa chất tẩy rửa mạnh.", "font_name": "Calibri", "size": 10.5, "color": c_charcoal, "space_after": 4},
        {"text": "• Kết hợp đầu nối KSD-D và KSM-D đạt tiêu chuẩn vệ sinh, không đọng nước.", "font_name": "Calibri", "size": 10.5, "color": c_charcoal, "space_after": 4},
        {"text": "• Cung cấp giải pháp bảo vệ cáp khép kín từ tủ điện đến động cơ máy.", "font_name": "Calibri", "size": 10.5, "color": c_charcoal, "space_after": 4}
    ]
    add_textbox_styled(slide7, x_right + Inches(0.2), y_split + Inches(2.6), w_split - Inches(0.4), Inches(2.3), conduit_text)

    # ==========================================
    # SLIDE 8: DÂY BUỘC CÁP PHÁT HIỆN ĐƯỢC (DETECTABLE)
    # ==========================================
    slide8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide8, c_grey_bg)
    add_slide_header(slide8, "07", "Dây Buộc Cáp & Nhãn Ghi Phát Hiện Được")
    
    # Left Column: Features & Safety
    left_x = Inches(0.6)
    left_y = Inches(1.6)
    left_w = Inches(6.0)
    
    left_text = [
        {"text": "GIẢI PHÁP PHÒNG CHỐNG DỊ VẬT THỰC PHẨM", "font_name": "Trebuchet MS", "size": 15, "bold": True, "color": c_navy, "space_after": 15},
        {"text": "• Nhựa Nylon tích hợp hạt kim loại phân bố đều:", "font_name": "Calibri", "size": 12, "bold": True, "color": c_blue, "space_after": 2},
        {"text": "Dây buộc cáp (cable ties) và thẻ nhãn ghi nhãn được pha trộn bột kim loại, giúp máy dò kim loại (metal detector) hoặc hệ thống máy quét tia X-ray trong dây chuyền thực phẩm dễ dàng phát hiện nếu bị đứt rơi vào sản phẩm.", "font_name": "Calibri", "size": 11, "color": c_charcoal, "space_after": 12},
        {"text": "• Màu xanh lam (Blue) đặc trưng nhận biết nhanh:", "font_name": "Calibri", "size": 12, "bold": True, "color": c_blue, "space_after": 2},
        {"text": "Giúp người vận hành dễ dàng nhận diện dị vật bằng mắt thường trong quá trình kiểm tra ngoại quan sản phẩm thực phẩm.", "font_name": "Calibri", "size": 11, "color": c_charcoal, "space_after": 12},
        {"text": "• Ứng dụng dán nhãn vệ sinh:", "font_name": "Calibri", "size": 12, "bold": True, "color": c_blue, "space_after": 2},
        {"text": "Nhãn ghi bằng nhựa và thép không gỉ in bằng máy in chuyên dụng mp-LM 3, bền bỉ trước môi trường tẩy rửa độ ẩm cao.", "font_name": "Calibri", "size": 11, "color": c_charcoal, "space_after": 4}
    ]
    add_textbox_styled(slide8, left_x, left_y, left_w, Inches(5.0), left_text)
    
    # Right Column: Visual Product Image
    right_x = Inches(7.0)
    right_y = Inches(1.6)
    right_w = Inches(5.7)
    right_h = Inches(4.5)
    
    add_card(slide8, right_x, right_y, right_w, right_h)
    add_fit_picture_styled(slide8, os.path.join(prod_img_dir, "SUV_EW_R_PP_produkt_01.png"), right_x + Inches(0.2), right_y + Inches(0.2), right_w - Inches(0.4), right_h - Inches(0.4))
    
    # Small overlay card highlighting Detectable feature
    det_x = right_x + Inches(0.4)
    det_y = right_y + Inches(3.2)
    det_w = Inches(2.2)
    det_h = Inches(0.9)
    add_card(slide8, det_x, det_y, det_w, det_h, fill_color=c_cyan, border_color=c_cyan)
    p_det = slide8.shapes.add_textbox(det_x, det_y + Inches(0.05), det_w, det_h).text_frame.paragraphs[0]
    p_det.text = "DETECTABLE TIES"
    p_det.font.name = 'Plus Jakarta Sans'
    p_det.font.size = Pt(13)
    p_det.font.bold = True
    p_det.font.color.rgb = c_white
    p_det.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 9: DỰ ÁN THỰC TẾ & KHÁCH HÀNG TIÊU BIỂU
    # ==========================================
    slide9 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide9, c_grey_bg)
    add_slide_header(slide9, "08", "Dự Án Thực Tế & Khách Hàng Tiêu Biểu")
    
    # Case Study 1: Velec Systems
    x_cs1 = Inches(0.6)
    y_cs = Inches(1.6)
    w_cs = Inches(5.8)
    h_cs = Inches(5.2)
    
    add_card(slide9, x_cs1, y_cs, w_cs, h_cs)
    add_fit_picture_styled(slide9, os.path.join(real_img_dir, "Velec Systems (nhà sản xuất máy làm xúc xích)_1.png"), x_cs1 + Inches(0.2), y_cs + Inches(0.2), w_cs - Inches(0.4), Inches(3.0))
    
    cs1_text = [
        {"text": "VELEC SYSTEMS (PHÁP) - MÁY LÀM XÚC XÍCH", "font_name": "Trebuchet MS", "size": 12, "bold": True, "color": c_navy, "space_after": 8},
        {"text": "• Ứng dụng giải pháp tấm dẫn cáp hợp vệ sinh của Murrplastik trên thân máy chính.", "font_name": "Calibri", "size": 10.5, "color": c_charcoal, "space_after": 4},
        {"text": "• Giúp nhà máy rút ngắn đáng kể chu kỳ vệ sinh xịt rửa vòi phun áp lực cao.", "font_name": "Calibri", "size": 10.5, "color": c_charcoal, "space_after": 4},
        {"text": "• Đảm bảo tuyệt đối an toàn vệ sinh dịch tễ cho sản phẩm xúc xích đầu ra.", "font_name": "Calibri", "size": 10.5, "color": c_charcoal, "space_after": 4}
    ]
    add_textbox_styled(slide9, x_cs1 + Inches(0.2), y_cs + Inches(3.4), w_cs - Inches(0.4), Inches(1.6), cs1_text)
    
    # Case Study 2: Zuger AG
    x_cs2 = Inches(6.93)
    add_card(slide9, x_cs2, y_cs, w_cs, h_cs)
    add_fit_picture_styled(slide9, os.path.join(real_img_dir, "Züger AG (Thụy Sĩ) Các đầu nối cáp kim loại.png"), x_cs2 + Inches(0.2), y_cs + Inches(0.2), w_cs - Inches(0.4), Inches(3.0))
    
    cs2_text = [
        {"text": "ZÜGER AG (THỤY SĨ) - NHÀ MÁY CHẾ BIẾN SỮA", "font_name": "Trebuchet MS", "size": 12, "bold": True, "color": c_navy, "space_after": 8},
        {"text": "• Ứng dụng các đầu nối cáp kim loại và tấm dẫn cáp thép không gỉ tại các khu vực chế biến ướt.", "font_name": "Calibri", "size": 10.5, "color": c_charcoal, "space_after": 4},
        {"text": "• Chịu được áp lực nước xịt rửa mạnh kèm hóa chất tẩy rửa khắc nghiệt hàng ngày.", "font_name": "Calibri", "size": 10.5, "color": c_charcoal, "space_after": 4},
        {"text": "• Hệ thống hoạt động tin cậy nhiều năm mà không xảy ra sự cố rò rỉ hay ô nhiễm sinh học.", "font_name": "Calibri", "size": 10.5, "color": c_charcoal, "space_after": 4}
    ]
    add_textbox_styled(slide9, x_cs2 + Inches(0.2), y_cs + Inches(3.4), w_cs - Inches(0.4), Inches(1.6), cs2_text)

    # ==========================================
    # SLIDE 10: LIÊN HỆ & HỢP TÁC (CONTACT SLIDE)
    # ==========================================
    slide10 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide10, c_navy)
    
    # Left accent background
    left_accent10 = slide10.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(7.5), Inches(7.5)
    )
    left_accent10.fill.solid()
    left_accent10.fill.fore_color.rgb = RGBColor(2, 3, 70)
    left_accent10.line.fill.background()
    
    # Left Column: T&T Vina info
    info_text = [
        {"text": "CÔNG TY TNHH CÔNG NGHIỆP T&T VINA", "font_name": "Trebuchet MS", "size": 22, "bold": True, "color": c_white, "space_after": 15},
        {"text": "Đối tác đồng hành tin cậy cho giải pháp đi dây cáp cơ điện nhà máy.", "font_name": "Calibri", "size": 13, "color": RGBColor(200, 220, 240), "space_after": 25},
        {"text": "• HÀNG CHÍNH HÃNG:", "font_name": "Trebuchet MS", "size": 12, "bold": True, "color": c_cyan, "space_after": 2},
        {"text": "Nhập khẩu trực tiếp từ Murrplastik (Đức), đầy đủ CO/CQ và chứng chỉ chất lượng FDA/EHEDG/Ecolab.", "font_name": "Calibri", "size": 11, "color": c_white, "space_after": 12},
        {"text": "• HỖ TRỢ KỸ THUẬT:", "font_name": "Trebuchet MS", "size": 12, "bold": True, "color": c_cyan, "space_after": 2},
        {"text": "Tư vấn kỹ thuật tại công trình, bóc tách bản vẽ và thiết kế cá nhân hóa KDP theo yêu cầu nhà máy.", "font_name": "Calibri", "size": 11, "color": c_white, "space_after": 12},
        {"text": "• DỊCH VỤ NHANH CHÓNG:", "font_name": "Trebuchet MS", "size": 12, "bold": True, "color": c_cyan, "space_after": 2},
        {"text": "Thời gian giao hàng tối ưu, dịch vụ hậu mãi chu đáo cho mọi dự án.", "font_name": "Calibri", "size": 11, "color": c_white, "space_after": 4}
    ]
    add_textbox_styled(slide10, Inches(0.6), Inches(1.2), Inches(6.3), Inches(5.5), info_text)
    
    # Right Column: Contact details & QR code box
    right_x10 = Inches(8.2)
    right_y10 = Inches(1.2)
    right_w10 = Inches(4.5)
    
    contact_text = [
        {"text": "THÔNG TIN LIÊN HỆ CHÀO DỰ ÁN", "font_name": "Trebuchet MS", "size": 14, "bold": True, "color": c_cyan, "space_after": 20},
        {"text": "Hotline bán hàng:", "font_name": "Calibri", "size": 11, "bold": True, "color": RGBColor(200, 220, 240), "space_after": 2},
        {"text": "+84 973 363 824", "font_name": "Trebuchet MS", "size": 18, "bold": True, "color": c_white, "space_after": 15},
        {"text": "Email gửi yêu cầu báo giá:", "font_name": "Calibri", "size": 11, "bold": True, "color": RGBColor(200, 220, 240), "space_after": 2},
        {"text": "sales@murrplastik-vn.com", "font_name": "Trebuchet MS", "size": 14, "bold": True, "color": c_white, "space_after": 15},
        {"text": "Website chính thức tại Việt Nam:", "font_name": "Calibri", "size": 11, "bold": True, "color": RGBColor(200, 220, 240), "space_after": 2},
        {"text": "murrplastikvn.com", "font_name": "Trebuchet MS", "size": 14, "bold": True, "color": c_white, "space_after": 15}
    ]
    add_textbox_styled(slide10, right_x10, right_y10, right_w10, Inches(3.0), contact_text)
    
    # Stylized QR Code placeholder
    qr_x = right_x10
    qr_y = Inches(4.6)
    qr_w = Inches(1.8)
    qr_h = Inches(1.8)
    add_card(slide10, qr_x, qr_y, qr_w, qr_h, fill_color=c_white, border_color=c_cyan)
    
    # Inner QR detail lines representation
    qr_inner = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, qr_x + Inches(0.15), qr_y + Inches(0.15), qr_w - Inches(0.3), qr_h - Inches(0.3))
    qr_inner.fill.solid()
    qr_inner.fill.fore_color.rgb = c_navy
    qr_inner.line.fill.background()
    try:
        qr_inner.adjustments[0] = 0.05
    except:
        pass
    p_qr = qr_inner.text_frame.paragraphs[0]
    p_qr.text = "WEBSITE\nQR CODE"
    p_qr.font.name = 'Plus Jakarta Sans'
    p_qr.font.size = Pt(11)
    p_qr.font.bold = True
    p_qr.font.color.rgb = c_cyan
    p_qr.alignment = PP_ALIGN.CENTER
    
    # Save Presentation
    output_pptx_path = os.path.join(root_dir, "f&b_hygienic_cabling_murrplastik.pptx")
    try:
        prs.save(output_pptx_path)
        print(f"Presentation saved successfully to: {output_pptx_path}")
    except PermissionError:
        alternative_pptx_path = os.path.join(root_dir, "f&b_hygienic_cabling_murrplastik_updated.pptx")
        prs.save(alternative_pptx_path)
        print(f"Warning: {output_pptx_path} is locked/open. Saved alternative to: {alternative_pptx_path}")

if __name__ == "__main__":
    main()
